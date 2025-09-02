"""
Générateur de présentation PowerPoint à partir des analyses de prestations
"""

import os
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json

def charger_donnees_analyse():
    """Charge les données depuis le fichier source en utilisant la même fonction que le script principal"""
    try:
        # Importer et utiliser la fonction charger du script principal
        import sys
        sys.path.append('.')
        from analyse_prestations_full import charger
        return charger('Classeur1.xlsx')
    except Exception as e:
        print(f"Erreur lors du chargement des données: {e}")
        # Fallback vers la méthode manuelle
        try:
            df = pd.read_excel('Classeur1.xlsx')
            # Appliquer les mêmes filtres que le script principal
            date_col = next((c for c in df.columns if c.lower() == 'date'), None)
            if date_col:
                df = df.copy()
                df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
                df = df.dropna(subset=[date_col])
                df = df[df[date_col].dt.year.isin([2024, 2025])]
            return df
        except Exception as e2:
            print(f"Erreur lors du fallback: {e2}")
            return None

def calculer_indicateurs_cles(df):
    """Calcule les indicateurs clés en utilisant les mêmes fonctions que le script principal"""
    if df is None or df.empty:
        return {}
    
    try:
        # Importer et utiliser la fonction compute_global du script principal
        import sys
        sys.path.append('.')
        from analyse_prestations_full import compute_global
        indicateurs = compute_global(df)
        
        # Mapper les noms des indicateurs pour la compatibilité
        mapping = {
            'nb_prestations': 'nombre_prestations',
            'nb_mutualistes_distincts': 'mutualistes_distincts',
            'taux_recours_pct': 'taux_recours',
            'taux_acceptation_pct': 'taux_acceptation'
        }
        
        for old_key, new_key in mapping.items():
            if old_key in indicateurs and new_key not in indicateurs:
                indicateurs[new_key] = indicateurs[old_key]
        
        return indicateurs
        
    except Exception as e:
        print(f"Erreur lors de l'import des fonctions principales: {e}")
        # Fallback vers le calcul manuel
        return calculer_indicateurs_fallback(df)

def calculer_indicateurs_fallback(df):
    """Calcul de fallback si l'import du script principal échoue"""
    montant_col = next((c for c in df.columns if c.lower().startswith('montant')), None)
    adherent_col = next((c for c in df.columns if 'adherent_code' in c.lower()), None)
    
    indicateurs = {}
    
    if montant_col:
        indicateurs['montant_total'] = df[montant_col].sum()
        indicateurs['montant_moyen'] = df[montant_col].mean()
    
    indicateurs['nombre_prestations'] = len(df)
    
    if adherent_col:
        indicateurs['mutualistes_distincts'] = df[adherent_col].nunique()
        indicateurs['taux_recours'] = (df[adherent_col].nunique() / 5284) * 100  # 5284 = total mutualistes éligibles

    if 'statut' in df.columns:
        total = len(df)
        acceptes = (df['statut'] == 'accepté').sum() + (df['statut'] == 'accepte').sum()
        indicateurs['taux_acceptation'] = (acceptes / total * 100) if total else 0
    
    return indicateurs

def ajouter_slide_titre(prs):
    """Ajoute la slide de titre"""
    slide_layout = prs.slide_layouts[0]  # Titre
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Analyse des Prestations MUPOL"
    subtitle.text = "Rapport Analytique 2024-2025\nTableau de Bord Exécutif"
    
    # Style du titre
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def ajouter_slide_indicateurs_cles(prs, indicateurs):
    """Ajoute la slide des indicateurs clés"""
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Indicateurs Clés de Performance"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Création des boîtes d'indicateurs
    boites = [
        ("Montant Total", f"{indicateurs.get('montant_total', 0):,.0f} FCFA", RGBColor(52, 152, 219)),
        ("Prestations", f"{indicateurs.get('nombre_prestations', 0):,}", RGBColor(46, 204, 113)),
        ("Mutualistes", f"{indicateurs.get('mutualistes_distincts', 0):,}", RGBColor(155, 89, 182)),
        ("Taux de Recours", f"{indicateurs.get('taux_recours', 0):.1f}%", RGBColor(230, 126, 34)),
        ("Montant Moyen", f"{indicateurs.get('montant_moyen', 0):,.0f} FCFA", RGBColor(231, 76, 60)),
        ("Taux d'Acceptation", f"{indicateurs.get('taux_acceptation', 0):.1f}%", RGBColor(26, 188, 156))
    ]
    
    # Disposition en grille 2x3
    for i, (titre, valeur, couleur) in enumerate(boites):
        row = i // 3
        col = i % 3
        
        x = Inches(0.5 + col * 3.2)
        y = Inches(2 + row * 2.5)
        width = Inches(3)
        height = Inches(2)
        
        # Boîte colorée
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = couleur
        shape.line.color.rgb = couleur
        
        # Texte du titre
        text_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.6))
        text_frame = text_box.text_frame
        text_para = text_frame.paragraphs[0]
        text_para.text = titre
        text_para.font.size = Pt(14)
        text_para.font.bold = True
        text_para.font.color.rgb = RGBColor(255, 255, 255)
        text_para.alignment = PP_ALIGN.CENTER
        
        # Texte de la valeur
        value_box = slide.shapes.add_textbox(x, y + Inches(0.8), width, Inches(1))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = valeur
        value_para.font.size = Pt(18)
        value_para.font.bold = True
        value_para.font.color.rgb = RGBColor(255, 255, 255)
        value_para.alignment = PP_ALIGN.CENTER

def ajouter_slide_graphique(prs, image_path, titre, description=""):
    """Ajoute une slide avec un graphique"""
    if not os.path.exists(image_path):
        return
    
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = titre
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Image
    try:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(9), height=Inches(5.5))
    except Exception as e:
        print(f"Erreur lors de l'ajout de l'image {image_path}: {e}")
    
    # Description si fournie
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = description
        desc_para.font.size = Pt(12)
        desc_para.alignment = PP_ALIGN.CENTER

def ajouter_slide_graphique_avec_tableau(prs, image_path, titre, df, type_analyse):
    """Ajoute une slide avec un tableau en haut et le graphique en dessous"""
    if not os.path.exists(image_path):
        return
    
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = titre
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Tableau en haut
    try:
        table_data = generer_donnees_tableau(df, type_analyse)
        if table_data is not None and len(table_data) > 0:
            ajouter_tableau_donnees(slide, table_data, x=Inches(0.5), y=Inches(0.8), width=Inches(9), height=Inches(2.5))
    except Exception as e:
        print(f"Erreur lors de l'ajout du tableau pour {type_analyse}: {e}")
    
    # Graphique en bas
    try:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(3.5))
    except Exception as e:
        print(f"Erreur lors de l'ajout de l'image {image_path}: {e}")
        return

def generer_donnees_tableau(df, type_analyse):
    """Génère les données de tableau selon le type d'analyse"""
    try:
        montant_col = next((c for c in df.columns if c.lower().startswith('montant')), None)
        if not montant_col:
            return None
            
        if type_analyse == "repartition_actes":
            if 'acte' in df.columns:
                result = df.groupby('acte')[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Type d\'Acte', 'Montant (FCFA)', 'Nombre']
                result = result.sort_values('Montant (FCFA)', ascending=False)
                total_montant = result['Montant (FCFA)'].sum()
                result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result.head(8)
                
        elif type_analyse == "top_centres_montant":
            centre_col = next((c for c in df.columns if 'centre' in c.lower()), None)
            if centre_col:
                result = df.groupby(centre_col)[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Centre', 'Montant (FCFA)', 'Prestations']
                result = result.sort_values('Montant (FCFA)', ascending=False)
                return result.head(10)
                
        elif type_analyse == "top_centres_nombre":
            centre_col = next((c for c in df.columns if 'centre' in c.lower()), None)
            if centre_col:
                result = df.groupby(centre_col)[montant_col].agg(['count', 'sum']).reset_index()
                result.columns = ['Centre', 'Prestations', 'Montant (FCFA)']
                result = result.sort_values('Prestations', ascending=False)
                return result.head(10)
                
        elif type_analyse == "repartition_statut":
            if 'statut' in df.columns:
                result = df.groupby('statut')[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Statut', 'Montant (FCFA)', 'Nombre']
                total_montant = result['Montant (FCFA)'].sum()
                result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result
                
        elif type_analyse == "repartition_beneficiaires":
            if 'beneficiaire' in df.columns:
                result = df.groupby('beneficiaire')[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Bénéficiaire', 'Montant (FCFA)', 'Nombre']
                # Filtrer pour exclure "Mutualiste" et lignes de proportion
                result = result[~result['Bénéficiaire'].str.contains('Mutualiste|Proportion', case=False, na=False)]
                if len(result) > 0:
                    total_montant = result['Montant (FCFA)'].sum()
                    result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result
                
        elif type_analyse == "repartition_genre":
            genre_col = next((c for c in df.columns if 'genre' in c.lower()), None)
            if genre_col:
                result = df.groupby(genre_col)[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Genre', 'Montant (FCFA)', 'Nombre']
                total_montant = result['Montant (FCFA)'].sum()
                result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result
                
        elif type_analyse == "repartition_types":
            if 'type' in df.columns:
                result = df.groupby('type')[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Type', 'Montant (FCFA)', 'Nombre']
                result = result.sort_values('Montant (FCFA)', ascending=False)
                total_montant = result['Montant (FCFA)'].sum()
                result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result.head(8)
                
        elif type_analyse == "repartition_region":
            if 'region' in df.columns:
                result = df.groupby('region')[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Région', 'Montant (FCFA)', 'Nombre']
                result = result.sort_values('Montant (FCFA)', ascending=False)
                total_montant = result['Montant (FCFA)'].sum()
                result['Pourcentage'] = (result['Montant (FCFA)'] / total_montant * 100).round(1)
                return result.head(8)
                
        elif type_analyse == "evolution_mensuelle":
            if 'date' in df.columns:
                df_copy = df.copy()
                # Gestion améliorée des dates
                df_copy['date'] = pd.to_datetime(df_copy['date'], errors='coerce')
                df_copy = df_copy.dropna(subset=['date'])
                if len(df_copy) > 0:
                    df_copy['Mois'] = df_copy['date'].dt.to_period('M')
                    result = df_copy.groupby('Mois')[montant_col].agg(['sum', 'count']).reset_index()
                    result.columns = ['Mois', 'Montant (FCFA)', 'Prestations']
                    result['Mois'] = result['Mois'].astype(str)
                    return result.tail(12)  # 12 derniers mois
                
        elif type_analyse == "top_adherents":
            adherent_col = next((c for c in df.columns if 'adherent' in c.lower() and 'code' in c.lower()), None)
            if adherent_col:
                result = df.groupby(adherent_col)[montant_col].agg(['sum', 'count']).reset_index()
                result.columns = ['Code Adhérent', 'Montant (FCFA)', 'Prestations']
                result = result.sort_values('Montant (FCFA)', ascending=False)
                return result.head(10)
                
        elif type_analyse == "evolution_actes_2024":
            if 'date' in df.columns and 'acte' in df.columns:
                df_copy = df.copy()
                df_copy['date'] = pd.to_datetime(df_copy['date'], errors='coerce')
                df_copy = df_copy.dropna(subset=['date'])
                # Filtrer pour 2024 seulement
                df_2024 = df_copy[df_copy['date'].dt.year == 2024]
                if len(df_2024) > 0:
                    df_2024['Mois'] = df_2024['date'].dt.to_period('M')
                    # Grouper par mois et acte
                    result = df_2024.groupby(['Mois', 'acte'])[montant_col].agg(['sum', 'count']).reset_index()
                    result.columns = ['Mois', 'Acte', 'Montant (FCFA)', 'Prestations']
                    result['Mois'] = result['Mois'].astype(str)
                    # Prendre les top actes
                    top_actes = df_2024.groupby('acte')[montant_col].sum().nlargest(5).index.tolist()
                    result = result[result['Acte'].isin(top_actes)]
                    return result.sort_values(['Mois', 'Montant (FCFA)'], ascending=[True, False])
                    
        elif type_analyse == "evolution_actes_2025":
            if 'date' in df.columns and 'acte' in df.columns:
                df_copy = df.copy()
                df_copy['date'] = pd.to_datetime(df_copy['date'], errors='coerce')
                df_copy = df_copy.dropna(subset=['date'])
                # Filtrer pour 2025 seulement
                df_2025 = df_copy[df_copy['date'].dt.year == 2025]
                if len(df_2025) > 0:
                    df_2025['Mois'] = df_2025['date'].dt.to_period('M')
                    # Grouper par mois et acte
                    result = df_2025.groupby(['Mois', 'acte'])[montant_col].agg(['sum', 'count']).reset_index()
                    result.columns = ['Mois', 'Acte', 'Montant (FCFA)', 'Prestations']
                    result['Mois'] = result['Mois'].astype(str)
                    # Prendre les top actes
                    top_actes = df_2025.groupby('acte')[montant_col].sum().nlargest(5).index.tolist()
                    result = result[result['Acte'].isin(top_actes)]
                    return result.sort_values(['Mois', 'Montant (FCFA)'], ascending=[True, False])
                
        return None
    
    except Exception as e:
        print(f"Erreur lors de la génération des données pour {type_analyse}: {e}")
        return None

def ajouter_tableau_donnees(slide, data, x, y, width, height):
    """Ajoute un tableau de données à la slide"""
    if data is None or len(data) == 0:
        return
    
    # Limiter le nombre de lignes si nécessaire
    max_rows = min(12, len(data))  # Maximum 12 lignes + en-tête
    
    rows, cols = max_rows + 1, len(data.columns)
    table = slide.shapes.add_table(rows, cols, x, y, width, height).table
    
    # En-têtes
    for i, col_name in enumerate(data.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Données
    for i in range(min(max_rows, len(data))):
        for j, col_name in enumerate(data.columns):
            cell = table.cell(i + 1, j)
            value = data.iloc[i, j]
            
            # Formatage selon le type de données
            if 'Montant' in col_name and 'FCFA' in col_name:
                cell.text = f"{value:,.0f}"
            elif 'Pourcentage' in col_name or col_name == '%':
                cell.text = f"{value:.1f}%"
            elif isinstance(value, (int, float)) and value > 1000:
                cell.text = f"{value:,.0f}"
            else:
                # Tronquer les textes longs
                text = str(value)
                cell.text = text[:30] + "..." if len(text) > 30 else text
            
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            
            # Couleur alternée
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)

def ajouter_slide_tableau_top_centres(prs, df):
    """Ajoute une slide avec le tableau des top centres"""
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Top 10 Centres de Santé"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    try:
        # Calcul des données pour le tableau
        centre_col = next((c for c in df.columns if 'centre' in c.lower()), None)
        montant_col = next((c for c in df.columns if c.lower().startswith('montant')), None)
        
        if centre_col and montant_col:
            top_centres = df.groupby(centre_col)[montant_col].agg(['sum', 'count']).reset_index()
            top_centres.columns = ['Centre', 'Montant Total', 'Nombre Prestations']
            top_centres = top_centres.sort_values('Montant Total', ascending=False).head(10)
            
            # Création du tableau
            rows, cols = min(11, len(top_centres) + 1), 3  # +1 pour l'en-tête
            x, y, width, height = Inches(1), Inches(1.2), Inches(8), Inches(5)
            
            table = slide.shapes.add_table(rows, cols, x, y, width, height).table
            
            # En-têtes
            headers = ['Centre de Santé', 'Montant (FCFA)', 'Prestations']
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # Données
            for i, (_, row) in enumerate(top_centres.iterrows()):
                if i >= 10:  # Limite à 10 lignes
                    break
                table.cell(i + 1, 0).text = str(row['Centre'])[:40]  # Tronquer si trop long
                table.cell(i + 1, 1).text = f"{row['Montant Total']:,.0f}"
                table.cell(i + 1, 2).text = f"{row['Nombre Prestations']:,}"
                
                # Style des cellules
                for j in range(3):
                    cell = table.cell(i + 1, j)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    except Exception as e:
        print(f"Erreur lors de la création du tableau des centres: {e}")

def ajouter_slide_repartition_actes(prs, df):
    """Ajoute une slide avec la répartition par actes"""
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Répartition par Type d'Acte"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    try:
        montant_col = next((c for c in df.columns if c.lower().startswith('montant')), None)
        
        if montant_col and 'acte' in df.columns:
            repartition = df.groupby('acte')[montant_col].agg(['sum', 'count']).reset_index()
            repartition.columns = ['Type d\'Acte', 'Montant Total', 'Nombre']
            repartition = repartition.sort_values('Montant Total', ascending=False)
            repartition['Pourcentage'] = (repartition['Montant Total'] / repartition['Montant Total'].sum() * 100)
            
            # Création du tableau
            rows, cols = min(11, len(repartition) + 1), 4
            x, y, width, height = Inches(1), Inches(1.2), Inches(8), Inches(5)
            
            table = slide.shapes.add_table(rows, cols, x, y, width, height).table
            
            # En-têtes
            headers = ['Type d\'Acte', 'Montant (FCFA)', 'Nombre', '%']
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(46, 204, 113)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # Données
            for i, (_, row) in enumerate(repartition.iterrows()):
                if i >= 10:
                    break
                table.cell(i + 1, 0).text = str(row['Type d\'Acte']).replace('_', ' ').title()
                table.cell(i + 1, 1).text = f"{row['Montant Total']:,.0f}"
                table.cell(i + 1, 2).text = f"{row['Nombre']:,}"
                table.cell(i + 1, 3).text = f"{row['Pourcentage']:.1f}%"
                
                # Style
                for j in range(4):
                    cell = table.cell(i + 1, j)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 255, 240)
    
    except Exception as e:
        print(f"Erreur lors de la création de la répartition par actes: {e}")

def ajouter_slide_conclusion(prs, indicateurs):
    """Ajoute la slide de conclusion"""
    slide_layout = prs.slide_layouts[6]  # Vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Points Clés & Recommandations"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Points clés
    points = [
        f"✓ Volume d'activité: {indicateurs.get('nombre_prestations', 0):,} prestations traitées",
        f"✓ Taux de recours: {indicateurs.get('taux_recours', 0):.1f}% des mutualistes ont utilisé leurs prestations",
        f"✓ Montant moyen par prestation: {indicateurs.get('montant_moyen', 0):,.0f} FCFA",
        f"✓ Taux d'acceptation: {indicateurs.get('taux_acceptation', 0):.1f}% des demandes acceptées",
        "",
        "Recommandations:",
        "• Optimiser la communication sur l'utilisation des prestations",
        "• Renforcer le contrôle qualité des centres partenaires",
        "• Développer la prévention pour réduire les coûts moyens",
        "• Améliorer l'accès aux soins dans les régions sous-représentées"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    
    for i, point in enumerate(points):
        if i > 0:
            content_frame.add_paragraph()
        para = content_frame.paragraphs[i]
        para.text = point
        para.font.size = Pt(16) if point.startswith('✓') or point.startswith('•') else Pt(18)
        para.font.bold = point == "Recommandations:"
        para.font.color.rgb = RGBColor(0, 51, 102) if point.startswith('✓') else RGBColor(51, 51, 51)

def generer_presentation_powerpoint():
    """Génère la présentation PowerPoint complète"""
    print("Génération de la présentation PowerPoint...")
    
    # Chargement des données
    df = charger_donnees_analyse()
    if df is None:
        print("Erreur: Impossible de charger les données")
        return
    
    # Calcul des indicateurs
    indicateurs = calculer_indicateurs_cles(df)
    
    # Création de la présentation
    prs = Presentation()
    
    # Slides
    print("Ajout de la slide de titre...")
    ajouter_slide_titre(prs)
    
    print("Ajout des indicateurs clés...")
    ajouter_slide_indicateurs_cles(prs, indicateurs)
    
    print("Ajout des graphiques avec tableaux de données...")
    # Tous les graphiques avec leurs tableaux de données correspondants
    graphiques_avec_donnees = [
        ("figures/montant_par_acte.png", "Répartition des Montants par Type d'Acte", "repartition_actes"),
        ("figures/montant_par_acte_pie.png", "Répartition des Montants par Acte (Camembert)", "repartition_actes"),
        ("figures/top10_centres.png", "Top 10 Centres - Montants", "top_centres_montant"),
        ("figures/top10_centres_nbp.png", "Top 10 Centres - Nombre de Prestations", "top_centres_nombre"),
        ("figures/repartition_statut_pie.png", "Répartition par Statut", "repartition_statut"),
        ("figures/repartition_beneficiaire_pie.png", "Répartition par Bénéficiaires", "repartition_beneficiaires"),
        ("figures/repartition_genre_pie.png", "Répartition par Genre", "repartition_genre"),
        ("figures/montant_par_type.png", "Répartition par Type de Prestation", "repartition_types"),
        ("figures/repartition_region_montant_nombre.png", "Répartition par Région", "repartition_region"),
        ("figures/top10_adherents_montant.png", "Top 10 Adhérents par Montant", "top_adherents"),
        ("figures/evolution_mensuelle.png", "Évolution Mensuelle des Prestations", "evolution_mensuelle"),
        ("figures/montant_moyen_par_acte_comparison.png", "Comparaison Montants Moyens 2024 vs 2025", None),
        ("figures/evolution_trimestrielle.png", "Évolution Trimestrielle", None),
        ("figures/proportions_prestations_sans_pharmacie.png", "Proportions Prestations (hors Pharmacie)", None)
    ]
    
    for image_path, titre, type_analyse in graphiques_avec_donnees:
        if os.path.exists(image_path):
            if type_analyse:
                ajouter_slide_graphique_avec_tableau(prs, image_path, titre, df, type_analyse)
            else:
                ajouter_slide_graphique(prs, image_path, titre)
    
    # Graphiques supplémentaires sans tableaux spécifiques
    autres_graphiques = [
        ("figures/montant_par_sous_type.png", "Répartition par Sous-Type"),
        ("figures/repartition_province_montant_nombre.png", "Répartition par Province"),
        ("figures/evolution_mensuelle_actes_2024.png", "Évolution Mensuelle par Acte 2024", "evolution_actes_2024"),
        ("figures/evolution_mensuelle_actes_2025.png", "Évolution Mensuelle par Acte 2025", "evolution_actes_2025"),
        ("figures/pareto_actes.png", "Analyse Pareto des Actes"),
        ("figures/pareto_types.png", "Analyse Pareto des Types"),
        ("figures/scatter_types.png", "Analyse de Corrélation par Types")
    ]
    
    for item in autres_graphiques:
        if len(item) == 3:  # Graphique avec tableau
            image_path, titre, type_analyse = item
            if os.path.exists(image_path):
                ajouter_slide_graphique_avec_tableau(prs, image_path, titre, df, type_analyse)
        else:  # Graphique simple
            image_path, titre = item
            if os.path.exists(image_path):
                ajouter_slide_graphique(prs, image_path, titre)
    
    print("Ajout de la conclusion...")
    ajouter_slide_conclusion(prs, indicateurs)
    
    # Sauvegarde
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Presentation_Analyse_Prestations_MUPOL_{timestamp}.pptx"
    prs.save(filename)
    print(f"Présentation sauvegardée: {filename}")
    
    return filename

if __name__ == "__main__":
    try:
        generer_presentation_powerpoint()
        print("✓ Présentation PowerPoint générée avec succès!")
    except Exception as e:
        print(f"Erreur lors de la génération: {e}")
        import traceback
        traceback.print_exc()
